"""SVG → PPTX / PDF / 单文件 HTML 导出。

PPTX 路线（核心）：
  - 用 python-pptx 创建 16:9 (13.333 × 7.5 inch) 演示文稿
  - 每页一张幻灯片，铺满 add_picture(PNG fallback)
  - 用 lxml 注入 OOXML asvg:svgBlip 扩展，让 PowerPoint 2019+/365 显示矢量 SVG（可编辑）
  - 旧版 PowerPoint / WPS 显示 PNG fallback

依赖 shoot 命令产出的 shots/*.png。如缺失则报错让用户先跑 ppt shoot。

PDF 路线：把 deck.html 用 chrome --print-to-pdf 转为 PDF。
HTML 路线：deck.html 的单文件版本（SVG 内嵌为 data URI 而非外链）。
"""

from __future__ import annotations

import base64
import json
import re
import subprocess
import sys
import tempfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.parts.image import ImagePart
from pptx.util import Inches

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"


def _embed_fonts(pptx_path: Path, fonts: dict[str, Path]) -> None:
    """将字体文件嵌入 PPTX（OOXML post-processing）。fonts = {font_name: ttf_path}。"""
    import zipfile
    from io import BytesIO

    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    FONT_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"

    CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

    buf = BytesIO()
    with zipfile.ZipFile(pptx_path, "r") as zf, zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as out:
        names = zf.namelist()
        skip = {"ppt/presentation.xml", "ppt/_rels/presentation.xml.rels", "[Content_Types].xml"}
        for name in names:
            if name not in skip:
                out.writestr(name, zf.read(name))

        pres_xml = etree.fromstring(zf.read("ppt/presentation.xml"))
        rels_xml = etree.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
        ct_xml = etree.fromstring(zf.read("[Content_Types].xml"))

        efl = pres_xml.find(f"{{{P_NS}}}embeddedFontLst")
        if efl is None:
            efl = etree.SubElement(pres_xml, f"{{{P_NS}}}embeddedFontLst")

            # Count existing rIds
            existing_rids = {r.get("Id", "") for r in rels_xml}

            next_rid = 1
            while f"rId_font_{next_rid}" in existing_rids:
                next_rid += 1

            for font_name, font_path in fonts.items():
                if not font_path.exists():
                    continue
                rid = f"rId_font_{next_rid}"
                font_data = font_path.read_bytes()
                font_entry = f"ppt/fonts/font{next_rid}.ttf"
                out.writestr(font_entry, font_data)

                # Add relationship
                rel = etree.SubElement(rels_xml, f"{{{R_NS}}}Relationship")
                rel.set("Id", rid)
                rel.set("Type", FONT_REL)
                rel.set("Target", f"fonts/font{next_rid}.ttf")

                # Add embeddedFont entry
                ef = etree.SubElement(efl, f"{{{P_NS}}}embeddedFont")
                fn = etree.SubElement(ef, f"{{{P_NS}}}font")
                fn.set("typeface", font_name)
                reg = etree.SubElement(ef, f"{{{P_NS}}}regular")
                reg.set(f"{{{R_NS}}}id", rid)

                next_rid += 1

            # 注册 .ttf 到 [Content_Types].xml（缺失会导致 PowerPoint 提示修复）
            has_ttf = any(d.get("Extension") == "ttf" for d in ct_xml.findall(f"{{{CT_NS}}}Default"))
            if not has_ttf:
                ttf_default = etree.SubElement(ct_xml, f"{{{CT_NS}}}Default")
                ttf_default.set("Extension", "ttf")
                ttf_default.set("ContentType", "application/x-font-ttf")

            # Update zip entries
            out.writestr(
                "ppt/presentation.xml",
                etree.tostring(pres_xml, xml_declaration=True, encoding="UTF-8", standalone=True),
            )
            out.writestr(
                "[Content_Types].xml",
                etree.tostring(ct_xml, xml_declaration=True, encoding="UTF-8", standalone=True),
            )
            out.writestr(
                "ppt/_rels/presentation.xml.rels",
                etree.tostring(rels_xml, xml_declaration=True, encoding="UTF-8", standalone=True),
            )

    # Replace original
    pptx_path.write_bytes(buf.getvalue())


def to_pptx(ws: Path) -> dict:
    """默认 PPTX 路径：用 NativeRenderer 直接生成原生 shape，100% 可编辑。"""
    sys.path.insert(0, str(Path(__file__).parent))
    from native_render import render_pptx

    out_path = render_pptx(ws)

    # Embed fonts if available
    font_map = {
        "Noto Serif SC": Path("/tmp/NotoSerifSC.ttf"),
        "Noto Sans SC": Path("/tmp/NotoSansSC.ttf"),
        "IBM Plex Mono": Path("/tmp/IBMPlexMono.ttf"),
    }
    available = {k: v for k, v in font_map.items() if v.exists()}
    if available:
        _embed_fonts(out_path, available)
        print(f"  [embed] {len(available)} fonts embedded: {', '.join(available)}")

    return {"output": str(out_path)}


def to_pptx_svg(ws: Path) -> dict:
    """旧路径：把 shoot 截好的 PNG 嵌入 pptx 并注入 svgBlip 矢量引用。
    显示效果好（保留 SVG 渐变 / 光斑），但 PowerPoint 把整页当一个图片对象，
    需要右键"转换为形状"才能编辑文字。
    """
    slides_dir = ws / "slides"
    shots_dir = ws / "shots"
    if not slides_dir.is_dir():
        raise SystemExit(f"[export] 没找到 {slides_dir}，先跑 ppt scaffold")

    svgs = sorted(slides_dir.glob("*.svg"))
    if not svgs:
        raise SystemExit(f"[export] {slides_dir} 没有 SVG，先跑 ppt scaffold")

    missing_pngs = [s.stem for s in svgs if not (shots_dir / (s.stem + ".png")).exists()]
    if missing_pngs:
        raise SystemExit(
            f"[export] 缺少 PNG fallback：{missing_pngs}\n"
            f"  PPTX 需要 PNG 兜底（旧版 PowerPoint / WPS 不支持 SVG 时显示 PNG）。\n"
            f"  请先跑：ppt shoot {ws}"
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, svg_path in enumerate(svgs, start=1):
        png_path = shots_dir / (svg_path.stem + ".png")
        slide = prs.slides.add_slide(blank_layout)
        pic = slide.shapes.add_picture(
            str(png_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        _inject_svg_blip(slide, pic, svg_path, slide_idx=idx)

    out_path = ws / "deck-svg.pptx"
    prs.save(str(out_path))
    return {"output": str(out_path)}


def _inject_svg_blip(slide, picture, svg_path: Path, slide_idx: int) -> None:
    """把 SVG 文件作为新的 image part 加入 slide rels，
    并在 picture 的 a:blip 元素下注入 asvg:svgBlip 引用。

    PowerPoint 2019+ / Office 365 会优先渲染矢量 SVG，旧版本显示 PNG fallback。
    """
    svg_bytes = svg_path.read_bytes()
    partname = PackURI(f"/ppt/media/svg-deck-{slide_idx}.svg")
    img_part = ImagePart(partname, "image/svg+xml", slide.part.package, svg_bytes)
    rId = slide.part.relate_to(img_part, RT.IMAGE)

    blip = picture._element.find(f".//{{{A_NS}}}blip")
    if blip is None:
        return  # 理论不会，但兜底防呆
    extLst = etree.SubElement(blip, f"{{{A_NS}}}extLst")
    ext = etree.SubElement(extLst, f"{{{A_NS}}}ext")
    ext.set("uri", SVG_EXT_URI)
    svgBlip = etree.SubElement(ext, f"{{{ASVG_NS}}}svgBlip", nsmap={"asvg": ASVG_NS})
    svgBlip.set(f"{{{R_NS}}}embed", rId)


def to_pdf(ws: Path) -> dict:
    """用 chrome headless --print-to-pdf 把 deck.html 转 PDF。"""
    deck = ws / "deck.html"
    if not deck.exists():
        raise SystemExit(f"[export] 没找到 {deck}，先跑 ppt shoot 生成 deck.html")
    sys.path.insert(0, str(Path(__file__).parent))
    from shoot import find_chrome

    chrome = find_chrome()
    out_path = ws / "deck.pdf"
    with tempfile.TemporaryDirectory(prefix="ppt-export-") as tmpdir:
        cmd = [
            chrome,
            "--headless",
            "--disable-gpu",
            "--no-sandbox",
            "--no-first-run",
            "--no-default-browser-check",
            f"--user-data-dir={tmpdir}",
            f"--print-to-pdf={out_path}",
            "--print-to-pdf-no-header",
            f"file://{deck.resolve()}",
        ]
        subprocess.run(cmd, check=False, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=60)
    if not out_path.exists():
        raise SystemExit("[export] chrome --print-to-pdf 失败")
    return {"output": str(out_path)}


def to_html(ws: Path) -> dict:
    """生成单文件 deck-standalone.html：把 SVG 和 PNG 都内嵌为 data URI，可独立分发。"""
    deck = ws / "deck.html"
    if not deck.exists():
        raise SystemExit(f"[export] 没找到 {deck}，先跑 ppt shoot")

    slides_dir = ws / "slides"
    shots_dir = ws / "shots"
    svgs = sorted(slides_dir.glob("*.svg"))
    slides_inline = []
    for svg in svgs:
        svg_b64 = base64.b64encode(svg.read_bytes()).decode("ascii")
        png = shots_dir / (svg.stem + ".png")
        png_b64 = base64.b64encode(png.read_bytes()).decode("ascii") if png.exists() else ""
        slides_inline.append(
            {
                "svg": f"data:image/svg+xml;base64,{svg_b64}",
                "thumb": f"data:image/png;base64,{png_b64}",
                "name": svg.stem,
            }
        )

    html = deck.read_text(encoding="utf-8")
    html = re.sub(
        r"const SLIDES = .*?;",
        "const SLIDES = " + json.dumps(slides_inline, ensure_ascii=False) + ";",
        html,
        count=1,
    )
    out_path = ws / "deck-standalone.html"
    out_path.write_text(html, encoding="utf-8")
    return {"output": str(out_path)}
