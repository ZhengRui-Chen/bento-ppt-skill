"""Native PPTX renderer：把 layout.json 直接渲染为 PowerPoint 原生 shape。

布局 100% 可编辑：单击文字进入编辑、单击 shape 拖动调色，零需"转换为形状"操作。

跟 SVG renderer 共享 layout.json 和 theme manifest 契约。
SVG 路径仍然用于 HTML 预览 / shoot 截图 / PDF。
"""

from __future__ import annotations

import json
import os
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SKILL_DIR = Path(os.environ.get("CLAUDE_SKILL_DIR", str(Path.home() / ".claude/skills/ppt-agent"))).resolve()
if not SKILL_DIR.exists():
    SKILL_DIR = Path(__file__).resolve().parent.parent  # fallback: dev clone
SLIDE_W_INCH = 13.333
SLIDE_H_INCH = 7.5
VIEWPORT_W = 1280
VIEWPORT_H = 720

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _hex(s: str) -> RGBColor:
    """'#7c5cff' → RGBColor(0x7c, 0x5c, 0xff)。"""
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _px_to_pt(px: float) -> float:
    """SVG px 单位 → PowerPoint pt（1 pt = 1.333 px @ 96 DPI）。"""
    return px * 0.75


class NativeRenderer:
    """layout.json + theme manifest → 原生 shape pptx。

    SVG 坐标系 (1280×720) 直接映射到 16:9 13.333"×7.5" 幻灯片。
    所有 shape 使用绝对 EMU 坐标（无嵌套坐标系）。
    """

    def __init__(self, theme_name: str = "bento-tech"):
        self.theme_name = theme_name
        self.theme = self._load_manifest(theme_name)
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W_INCH)
        self.prs.slide_height = Inches(SLIDE_H_INCH)
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        if slide_width is None or slide_height is None:
            raise RuntimeError("presentation slide size was not initialized")
        # 每 SVG 单位对应多少 EMU
        self.x_unit = slide_width / VIEWPORT_W
        self.y_unit = slide_height / VIEWPORT_H

    def _load_manifest(self, name: str) -> dict:
        path = SKILL_DIR / "themes" / name / "manifest.json"
        if not path.exists():
            raise SystemExit(f"[native_render] theme not found: {path}")
        return json.loads(path.read_text(encoding="utf-8"))

    @property
    def _is_light(self) -> bool:
        """Detect light themes so we can pick contrasting text colors on accent fills."""
        return self.theme["colors"]["bg_start"].lstrip("#").startswith(("F", "f"))

    # ---------- 坐标转换 ----------

    def _x(self, svg_x: float) -> Emu:
        return Emu(int(svg_x * self.x_unit))

    def _y(self, svg_y: float) -> Emu:
        return Emu(int(svg_y * self.y_unit))

    # ---------- 顶层 ----------

    def render_deck(self, layout: dict, out_path: Path) -> Path:
        pages = layout.get("pages", [])
        meta = layout.get("meta", {})
        for page in pages:
            self._render_slide(page, total=len(pages), meta=meta)
        out_path = Path(out_path)
        self.prs.save(str(out_path))
        return out_path

    def _render_slide(self, page: dict, total: int, meta: dict) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        self._add_background(slide)
        layout_name = page.get("layout")
        slot_defs = self.theme["layouts"].get(layout_name, {}).get("slots", {})
        pad = self.theme["spacing"]["card_padding"]
        for card in page.get("cards", []):
            slot = slot_defs.get(card.get("slot"))
            if not slot:
                continue
            self._add_card_bg(slide, slot)
            comp = card.get("component", "").replace("-", "_")
            handler = getattr(self, f"_render_{comp}", None)
            if handler is None:
                continue
            inner = {
                "x": slot["x"] + pad,
                "y": slot["y"] + pad,
                "w": slot["w"] - pad * 2,
                "h": slot["h"] - pad * 2,
            }
            handler(slide, inner, card.get("data", {}))
        self._add_footer(slide, meta, page, total)

    # ---------- 背景 + 页脚 ----------

    def _add_background(self, slide) -> None:
        """全屏底层矩形作为背景（不依赖 slide.background，那个被 master 继承覆盖）。
        放在所有 shape 之前画，确保 z-order 在最底。"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(0),
            Emu(0),
            self.prs.slide_width,
            self.prs.slide_height,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hex(self.theme["colors"]["bg_start"])
        bg.line.fill.background()
        # 锁住，防止用户误选误移（PowerPoint 端右键能解锁）
        # python-pptx 不直接支持 lockAspectRatio + selection lock，跳过

    def _add_footer(self, slide, meta: dict, page: dict, total: int) -> None:
        title = meta.get("title", "")
        page_n = page.get("page", 1)
        small_size = self.theme["type_scale"]["small"]
        # 左下：deck title
        self._add_textbox(
            slide,
            title,
            56,
            VIEWPORT_H - 26,
            600,
            20,
            font_size=small_size,
            color=self.theme["colors"]["text_muted"],
        )
        # 右下：页码
        self._add_textbox(
            slide,
            f"{page_n:02d} / {total:02d}",
            VIEWPORT_W - 200 - 56,
            VIEWPORT_H - 26,
            200,
            20,
            font_size=small_size,
            color=self.theme["colors"]["text_muted"],
            font_name="SF Mono",
            align="right",
        )

    # ---------- 卡片背景 ----------

    def _add_card_bg(self, slide, slot: dict) -> None:
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self._x(slot["x"]),
            self._y(slot["y"]),
            self._x(slot["w"]),
            self._y(slot["h"]),
        )
        # 卡片半透明白填充（玻璃感）
        self._set_solid_fill(
            sh,
            self.theme["colors"]["card_fill"],
            opacity=self.theme["colors"]["card_fill_opacity"],
        )
        # 描边带透明度
        self._set_line_alpha(
            sh,
            self.theme["colors"]["card_stroke"],
            opacity=self.theme["colors"]["card_stroke_opacity"],
            width_pt=self.theme["spacing"]["card_stroke_width"],
        )
        # 圆角调整：MSO ROUNDED_RECTANGLE adjustment 是相对短边的比例（0-0.5）
        target_radius_svg = self.theme["spacing"]["card_radius"]
        short_side = min(slot["w"], slot["h"])
        try:
            sh.adjustments[0] = min(0.5, target_radius_svg / short_side * 1.4)
        except (IndexError, AttributeError):
            pass

    # ---------- 共用工具：textbox ----------

    def _add_textbox(
        self,
        slide,
        text: str,
        svg_x: float,
        svg_y: float,
        svg_w: float,
        svg_h: float,
        *,
        font_size: float = 18,
        color: str = "#ffffff",
        bold: bool = False,
        italic: bool = False,
        align: str = "left",
        font_name: str | None = None,
        transparency: float = 0.0,
        letter_spacing: int = 0,
        word_wrap: bool = True,
        v_anchor: str = "top",
    ):
        """SVG 坐标 + 像素字号的 textbox。
        svg_y 是 textbox **顶部** y（非 baseline）；调用方负责把 SVG baseline 转 top。
        word_wrap=False 时数字/单行内容不会被 PowerPoint 自动换行（防溢出乱位置）。
        v_anchor='middle' 让文字在 textbox 内垂直居中（适合 quote 这种 wrap 后高度不定的场景）。
        """
        tb = slide.shapes.add_textbox(
            self._x(svg_x),
            self._y(svg_y),
            self._x(svg_w),
            self._y(svg_h),
        )
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = word_wrap
        if v_anchor == "middle":
            from pptx.enum.text import MSO_ANCHOR

            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        elif v_anchor == "bottom":
            from pptx.enum.text import MSO_ANCHOR

            tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0] if p.runs else None
        if run is not None:
            run.font.size = Pt(_px_to_pt(font_size))
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = _hex(color)
            run.font.name = font_name or "PingFang SC"
        if align == "right":
            p.alignment = PP_ALIGN.RIGHT
        elif align == "center":
            p.alignment = PP_ALIGN.CENTER
        else:
            p.alignment = PP_ALIGN.LEFT
        if (transparency > 0 or letter_spacing > 0) and run is not None:
            self._apply_rpr_extras(run, color, transparency, letter_spacing)
        return tb

    def _apply_rpr_extras(self, run, color: str, transparency: float, letter_spacing: int) -> None:
        """通过 lxml 给 run 的 a:rPr 加 transparency / character-spacing。"""
        rPr = run._r.get_or_add_rPr()
        if transparency > 0:
            for ch in list(rPr):
                if ch.tag == f"{{{A_NS}}}solidFill":
                    rPr.remove(ch)
            sf = etree.SubElement(rPr, f"{{{A_NS}}}solidFill")
            clr = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
            clr.set("val", color.lstrip("#"))
            alpha = etree.SubElement(clr, f"{{{A_NS}}}alpha")
            alpha.set("val", str(int((1.0 - transparency) * 100000)))
        if letter_spacing > 0:
            rPr.set("spc", str(letter_spacing))

    # ---------- 共用工具：badges ----------

    def _render_badges(self, slide, badges: list, svg_x: float, svg_y: float) -> None:
        """一排 badges 横向铺，22 高，间距 8。"""
        x_cursor = svg_x
        for b in badges:
            if isinstance(b, str):
                txt, var = b, "accent"
            else:
                txt = b.get("text", "")
                var = b.get("variant", "accent")
            w_est = sum(12 if "\u4e00" <= ch <= "\u9fff" else 7 for ch in txt) + 18
            if var == "success":
                bg = self.theme["colors"]["accent_success"]
                fg = "#ffffff"
            elif var == "warning":
                bg = self.theme["colors"]["accent_warning"]
                fg = "#ffffff"
            elif var == "muted":
                bg = "#ffffff"
                fg = self.theme["colors"]["text_secondary"]
            else:
                bg = self.theme["colors"]["accent_primary"]
                fg = "#ffffff"
            sh = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                self._x(x_cursor),
                self._y(svg_y),
                self._x(w_est),
                self._y(22),
            )
            self._set_solid_fill(sh, bg, opacity=0.22)
            self._set_line_alpha(sh, bg, opacity=0.55, width_pt=0.5)
            try:
                sh.adjustments[0] = 0.5
            except Exception:
                pass
            self._add_textbox(
                slide,
                txt,
                x_cursor,
                svg_y + 3,
                w_est,
                18,
                font_size=11,
                color=fg,
                bold=True,
                align="center",
                letter_spacing=50,
            )
            x_cursor += w_est + 8

    # ---------- 共用工具：fill alpha ----------

    def _set_solid_fill(self, shape, color_hex: str, opacity: float = 1.0) -> None:
        """实色填充 + 可选透明度。
        python-pptx 的 FillFormat 没有 transparency setter（曾是社区提案、未实现），
        必须用 lxml 给 <a:srgbClr> 加 <a:alpha val="..."/> 子元素。
        OOXML alpha val 是 0-100000 整数，100000 = 完全不透明。
        """
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex(color_hex)
        if opacity < 1.0:
            sp_pr = shape._element.spPr
            solid_fill = sp_pr.find(f"{{{A_NS}}}solidFill")
            if solid_fill is None:
                return
            clr = solid_fill[0]
            for child in list(clr):
                if child.tag == f"{{{A_NS}}}alpha":
                    clr.remove(child)
            alpha_el = etree.SubElement(clr, f"{{{A_NS}}}alpha")
            alpha_el.set("val", str(int(max(0.0, min(1.0, opacity)) * 100000)))

    def _set_line_alpha(self, shape, color_hex: str, opacity: float = 1.0, width_pt: float = 1) -> None:
        """描边色 + 可选透明度（同上原理，写到 line 的 solidFill）。"""
        shape.line.color.rgb = _hex(color_hex)
        shape.line.width = Pt(width_pt)
        if opacity < 1.0:
            sp_pr = shape._element.spPr
            ln = sp_pr.find(f"{{{A_NS}}}ln")
            if ln is None:
                return
            solid_fill = ln.find(f"{{{A_NS}}}solidFill")
            if solid_fill is None:
                return
            clr = solid_fill[0]
            for child in list(clr):
                if child.tag == f"{{{A_NS}}}alpha":
                    clr.remove(child)
            alpha_el = etree.SubElement(clr, f"{{{A_NS}}}alpha")
            alpha_el.set("val", str(int(max(0.0, min(1.0, opacity)) * 100000)))

    # ---------- 共用工具：进度条 + 装饰条 + accent 实色矩形 ----------

    def _add_solid_rect(
        self, slide, svg_x, svg_y, svg_w, svg_h, color: str, rounded: bool = False, transparency: float = 0.0
    ):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        sh = slide.shapes.add_shape(
            shape_type,
            self._x(svg_x),
            self._y(svg_y),
            self._x(svg_w),
            self._y(svg_h),
        )
        self._set_solid_fill(sh, color, opacity=1.0 - transparency)
        sh.line.fill.background()
        if rounded:
            try:
                sh.adjustments[0] = 0.5
            except Exception:
                pass
        return sh

    def _add_progress_bar(self, slide, svg_x, svg_y, svg_w, percent, label: str | None = None):
        # 底层灰色轨道
        track = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self._x(svg_x),
            self._y(svg_y),
            self._x(svg_w),
            self._y(6),
        )
        self._set_solid_fill(track, self.theme["colors"]["card_fill"], opacity=0.10)
        track.line.fill.background()
        try:
            track.adjustments[0] = 0.5
        except Exception:
            pass
        # 上层填充
        pct = max(0, min(100, int(percent)))
        if pct > 0:
            fill_w = int(svg_w * pct / 100)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                self._x(svg_x),
                self._y(svg_y),
                self._x(fill_w),
                self._y(6),
            )
            self._set_solid_fill(bar, self.theme["colors"]["accent_primary"], opacity=1.0)
            bar.line.fill.background()
            try:
                bar.adjustments[0] = 0.5
            except Exception:
                pass
        if label:
            self._add_textbox(
                slide,
                label,
                svg_x,
                svg_y + 14,
                svg_w,
                20,
                font_size=self.theme["type_scale"]["small"],
                color=self.theme["colors"]["text_secondary"],
            )

    # ---------- Component: card-stat ----------

    def _render_card_stat(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        # 字号自适应
        if W < 240:
            huge = 64
        elif W < 380:
            huge = 88
        else:
            huge = self.theme["type_scale"]["stat_huge"]
        unit_size = int(huge * 0.28)
        ch_size = 18 if W < 380 else 22
        ts = self.theme["type_scale"]
        v_str = str(data.get("value", ""))
        # value 文字宽估算
        vw = sum(huge if "\u4e00" <= ch <= "\u9fff" else int(huge * 0.55) for ch in v_str)

        y = 0
        if data.get("label"):
            y += ts["eyebrow"] + 4
            self._add_textbox(
                slide,
                data["label"],
                inner["x"],
                inner["y"] + y - ts["eyebrow"],
                W,
                ts["eyebrow"] + 6,
                font_size=ts["eyebrow"],
                color=self.theme["colors"]["accent_secondary"],
                bold=True,
                letter_spacing=300,
            )

        # 主数字
        y += 28 + huge
        self._add_textbox(
            slide,
            v_str,
            inner["x"],
            inner["y"] + y - huge,
            W,
            int(huge * 1.2),
            font_size=huge,
            color=self.theme["colors"]["text_primary"],
            bold=True,
            word_wrap=False,
        )

        unit = data.get("unit")
        unit_w_est = (len(unit) * unit_size * 0.6) if unit else 0
        if unit:
            if vw + 8 + unit_w_est <= W:
                # 同行
                self._add_textbox(
                    slide,
                    unit,
                    inner["x"] + vw + 8,
                    inner["y"] + y - unit_size - 2,
                    int(unit_w_est) + 20,
                    int(unit_size * 1.4),
                    font_size=unit_size,
                    color=self.theme["colors"]["text_secondary"],
                    bold=True,
                )
            else:
                # 换行到下方
                y += unit_size + 4
                self._add_textbox(
                    slide,
                    unit,
                    inner["x"],
                    inner["y"] + y - unit_size,
                    W,
                    int(unit_size * 1.4),
                    font_size=unit_size,
                    color=self.theme["colors"]["text_secondary"],
                    bold=True,
                )

        if data.get("sub_value"):
            sub_size = int(huge * 0.24)
            y += sub_size + 8
            self._add_textbox(
                slide,
                str(data["sub_value"]),
                inner["x"],
                inner["y"] + y - sub_size,
                W,
                int(sub_size * 1.4),
                font_size=sub_size,
                color=self.theme["colors"]["text_secondary"],
                bold=True,
            )

        if data.get("change"):
            cd = data.get("change_dir")
            if cd == "up":
                ch_color = self.theme["colors"]["accent_success"]
                arrow = "↑"
            elif cd == "down":
                ch_color = self.theme["colors"]["accent_warning"]
                arrow = "↓"
            else:
                ch_color = self.theme["colors"]["text_secondary"]
                arrow = "·"
            y += ch_size + 16
            self._add_textbox(
                slide,
                f"{arrow} {data['change']}",
                inner["x"],
                inner["y"] + y - ch_size,
                W,
                int(ch_size * 1.4),
                font_size=ch_size,
                color=ch_color,
                bold=True,
            )

        # 装饰条 + desc：动态 y
        deco_y = max(y + 24, H - 64)
        if deco_y < H - 12:
            self._add_solid_rect(
                slide,
                inner["x"],
                inner["y"] + deco_y,
                60,
                3,
                self.theme["colors"]["accent_primary"],
            )
            if data.get("desc") and deco_y + 30 < H:
                self._add_textbox(
                    slide,
                    data["desc"],
                    inner["x"],
                    inner["y"] + deco_y + 14,
                    W,
                    20,
                    font_size=ts["small"],
                    color=self.theme["colors"]["text_secondary"],
                )

    # ---------- Component: card-stack ----------

    def _render_card_stack(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        ts = self.theme["type_scale"]
        # 三档自适应
        if H >= 360:
            huge = 100 if W >= 600 else (80 if W >= 380 else 64)
            sec_gap, sec_label_gap, sec_value_size = 40, 30, ts["h3"]
        elif H >= 240:
            huge = 64 if W >= 380 else 48
            sec_gap, sec_label_gap, sec_value_size = 22, 22, 22
        else:
            huge = 48
            sec_gap, sec_label_gap, sec_value_size = 16, 18, 18
        unit_size = int(huge * 0.30)

        p = data.get("primary") or {}
        v_str = str(p.get("value", ""))
        vw = sum(huge if "\u4e00" <= ch <= "\u9fff" else int(huge * 0.55) for ch in v_str)

        # 上段
        y = 0
        if data.get("label"):
            y += ts["eyebrow"] + 4
            self._add_textbox(
                slide,
                data["label"],
                inner["x"],
                inner["y"] + y - ts["eyebrow"],
                W,
                ts["eyebrow"] + 6,
                font_size=ts["eyebrow"],
                color=self.theme["colors"]["accent_secondary"],
                bold=True,
                letter_spacing=300,
            )
        if data.get("badges"):
            y += 28
            self._render_badges(slide, data["badges"], inner["x"], inner["y"] + y - 18)
            y += 4
        top_h = y + 8

        # 下段：progress
        if data.get("progress"):
            bottom_h = 60 if data["progress"].get("label") else 30
        else:
            bottom_h = 0

        # 中段：内容居中
        has_sec = bool(data.get("secondary")) and H >= 180
        side_by_side = (H < 280) and (W >= 520) and has_sec
        primary_visual_h = int(huge * 0.85)

        if side_by_side or not has_sec:
            content_visual_h = primary_visual_h
        else:
            sec_visual_tail = int(sec_value_size * 0.2)
            content_visual_h = primary_visual_h + sec_gap + sec_label_gap + sec_visual_tail
        mid_avail = H - top_h - bottom_h
        mid_pad = max(0, (mid_avail - content_visual_h) // 2)
        primary_y = top_h + mid_pad + primary_visual_h

        # primary
        if p.get("value"):
            self._add_textbox(
                slide,
                v_str,
                inner["x"],
                inner["y"] + primary_y - huge,
                W,
                int(huge * 1.2),
                font_size=huge,
                color=self.theme["colors"]["text_primary"],
                bold=True,
                word_wrap=False,
            )
            if p.get("unit"):
                self._add_textbox(
                    slide,
                    p["unit"],
                    inner["x"] + vw + 8,
                    inner["y"] + primary_y - unit_size - 2,
                    int(unit_size * len(p["unit"]) * 0.7) + 20,
                    int(unit_size * 1.4),
                    font_size=unit_size,
                    color=self.theme["colors"]["accent_secondary"],
                    bold=True,
                )
            if p.get("suffix"):
                p_unit_w = int(len(p.get("unit", "")) * unit_size * 0.6)
                self._add_textbox(
                    slide,
                    p["suffix"],
                    inner["x"] + vw + 8 + p_unit_w + (12 if p.get("unit") else 0),
                    inner["y"] + primary_y - ts["body"] - 2,
                    int(W - vw - p_unit_w - 30),
                    int(ts["body"] * 1.4),
                    font_size=ts["body"],
                    color=self.theme["colors"]["text_secondary"],
                )

        # secondary
        if has_sec:
            sec = data["secondary"][:3]
            if side_by_side:
                p_unit_w = int(len(p.get("unit", "")) * unit_size * 0.6)
                p_suffix_w = int(len(p.get("suffix", "")) * ts["body"] * 0.55)
                primary_block_w = vw + 8 + p_unit_w + 12 + p_suffix_w + 24
                primary_block_w = max(primary_block_w, int(W * 0.5))
                sec_area_w = W - primary_block_w
                sec_col_w = sec_area_w // len(sec)
                sec_label_y = primary_y - int(primary_visual_h * 0.5)
                sec_value_y = sec_label_y + sec_label_gap
                for i, s in enumerate(sec):
                    sx = inner["x"] + primary_block_w + i * sec_col_w
                    self._add_textbox(
                        slide,
                        s.get("label", ""),
                        sx,
                        inner["y"] + sec_label_y - 12,
                        sec_col_w,
                        18,
                        font_size=12,
                        color=self.theme["colors"]["text_muted"],
                        letter_spacing=100,
                    )
                    self._add_textbox(
                        slide,
                        str(s.get("value", "")),
                        sx,
                        inner["y"] + sec_value_y - sec_value_size,
                        sec_col_w,
                        int(sec_value_size * 1.4),
                        font_size=sec_value_size,
                        color=self.theme["colors"]["text_primary"],
                        bold=True,
                    )
            else:
                col_w = W // len(sec)
                sec_label_y = primary_y + sec_gap
                sec_value_y = sec_label_y + sec_label_gap
                if sec_value_y < H - bottom_h - 2:
                    for i, s in enumerate(sec):
                        sx = inner["x"] + i * col_w
                        self._add_textbox(
                            slide,
                            s.get("label", ""),
                            sx,
                            inner["y"] + sec_label_y - 12,
                            col_w,
                            18,
                            font_size=12,
                            color=self.theme["colors"]["text_muted"],
                            letter_spacing=100,
                        )
                        self._add_textbox(
                            slide,
                            str(s.get("value", "")),
                            sx,
                            inner["y"] + sec_value_y - sec_value_size,
                            col_w,
                            int(sec_value_size * 1.4),
                            font_size=sec_value_size,
                            color=self.theme["colors"]["text_primary"],
                            bold=True,
                        )

        # 进度条
        if data.get("progress"):
            pr = data["progress"]
            bar_y = H - bottom_h
            self._add_progress_bar(
                slide,
                inner["x"],
                inner["y"] + bar_y,
                W,
                pr.get("percent", 0),
                pr.get("label"),
            )

    # ---------- Component: card-list ----------

    def _render_card_list(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        ts = self.theme["type_scale"]
        items_in = data.get("items") or []
        accent = data.get("accent", "default")
        if accent == "success":
            badge_fill = self.theme["colors"]["accent_success"]
            badge_text_color = "#ffffff" if self._is_light else "#0a3a1f"
        elif accent == "warning":
            badge_fill = self.theme["colors"]["accent_warning"]
            badge_text_color = "#ffffff" if self._is_light else "#3a2308"
        else:
            badge_fill = self.theme["colors"]["accent_primary"]
            badge_text_color = "#ffffff" if self._is_light else "#0a0e27"

        any_hl = any(isinstance(it, dict) and it.get("highlight") for it in items_in)

        y = 0
        if data.get("eyebrow"):
            y += ts["eyebrow"] + 4
            self._add_textbox(
                slide,
                data["eyebrow"],
                inner["x"],
                inner["y"] + y - ts["eyebrow"],
                W,
                ts["eyebrow"] + 6,
                font_size=ts["eyebrow"],
                color=self.theme["colors"]["accent_secondary"],
                bold=True,
                letter_spacing=300,
            )
        if data.get("title"):
            # 估算 title 是否会 wrap：CJK 字宽 ≈ font_size * 0.9，ASCII ≈ font_size * 0.5
            title = data["title"]
            title_w_est = sum(
                int(ts["h3"] * 0.9) if "\u4e00" <= ch <= "\u9fff" else int(ts["h3"] * 0.5) for ch in title
            )
            title_lines = 1 if title_w_est <= W - 16 else 2
            title_h = int(ts["h3"] * 1.25 * title_lines)
            self._add_textbox(
                slide,
                title,
                inner["x"],
                inner["y"] + y,
                W,
                title_h,
                font_size=ts["h3"],
                color=self.theme["colors"]["text_primary"],
                bold=True,
            )
            y += title_h + 8
        if not data.get("eyebrow") and not data.get("title"):
            y = 8
        list_top = y + 12

        has_desc = any(isinstance(it, dict) and it.get("desc") for it in items_in)
        min_row = 44 if has_desc else 32
        avail_h = H - list_top - 16
        max_fit = max(1, avail_h // min_row)
        # 选项策略：必须截断时优先保留 highlight 项；剩余按原顺序填补
        if len(items_in) <= max_fit:
            items_use = items_in
        else:
            hl_indices = [i for i, it in enumerate(items_in) if isinstance(it, dict) and it.get("highlight")]
            selected = set(hl_indices[:max_fit])
            for i in range(len(items_in)):
                if len(selected) >= max_fit:
                    break
                selected.add(i)
            items_use = [items_in[i] for i in sorted(selected)]
        n = len(items_use)
        row_h = max(min_row, min(80, avail_h // n if n else avail_h))

        for i, it in enumerate(items_use):
            ry = list_top + i * row_h
            if isinstance(it, str):
                it_title, it_desc, is_hl = it, None, False
            else:
                it_title = it.get("title", "")
                it_desc = it.get("desc")
                is_hl = bool(it.get("highlight"))
            # 高亮项左竖条：对齐"实际内容区"高度（序号方块顶到 desc 底），而不是含 gap 的整 row_h
            if is_hl:
                content_h = 50 if it_desc else 32
                self._add_solid_rect(
                    slide,
                    inner["x"] - 12,
                    inner["y"] + ry,
                    4,
                    content_h,
                    badge_fill,
                )
            # 序号方块
            badge_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                self._x(inner["x"]),
                self._y(inner["y"] + ry),
                self._x(32),
                self._y(32),
            )
            badge_box.fill.solid()
            badge_box.fill.fore_color.rgb = _hex(badge_fill)
            badge_box.line.fill.background()
            try:
                badge_box.adjustments[0] = 0.25
            except Exception:
                pass
            self._add_textbox(
                slide,
                f"{i + 1:02d}",
                inner["x"],
                inner["y"] + ry + 6,
                32,
                22,
                font_size=15,
                color=badge_text_color,
                bold=True,
                align="center",
            )
            # 标题
            title_color = (
                self.theme["colors"]["text_primary"]
                if (is_hl or not any_hl)
                else self.theme["colors"]["text_secondary"]
            )
            desc_color = (
                self.theme["colors"]["text_secondary"] if (is_hl or not any_hl) else self.theme["colors"]["text_muted"]
            )
            self._add_textbox(
                slide,
                it_title,
                inner["x"] + 48,
                inner["y"] + ry + 4,
                W - 48,
                24,
                font_size=18,
                color=title_color,
                bold=True,
            )
            if it_desc:
                self._add_textbox(
                    slide,
                    it_desc,
                    inner["x"] + 48,
                    inner["y"] + ry + 28,
                    W - 48,
                    22,
                    font_size=14,
                    color=desc_color,
                )

    # ---------- Component: card-quote ----------

    def _render_card_quote(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        ts = self.theme["type_scale"]
        is_banner = H < 280
        q_size = ts["h3"] if is_banner else ts["h2"]
        # 左侧渐变竖条（用纯色 accent）
        bar_h = H if is_banner else (H - (110 if data.get("author") else 30))
        self._add_solid_rect(
            slide,
            inner["x"],
            inner["y"],
            6,
            bar_h,
            self.theme["colors"]["accent_primary"],
            rounded=True,
        )

        quote = data.get("quote", "")

        if is_banner:
            # 横幅模式：quote 左侧 + author 右侧
            author_w = 300 if (data.get("author") or data.get("role")) else 0
            q_max_w = W - 32 - author_w - 24
            # textbox 高度 = H，用 vertical_anchor=middle 让 PowerPoint 自动垂直居中
            self._add_textbox(
                slide,
                quote,
                inner["x"] + 32,
                inner["y"],
                q_max_w,
                H,
                font_size=q_size,
                color=self.theme["colors"]["text_primary"],
                bold=True,
                v_anchor="middle",
            )
            if data.get("author") or data.get("role"):
                ax = inner["x"] + W - author_w + 16
                # author 区也用 vertical_anchor=middle
                a_lines = []
                if data.get("author"):
                    a_lines.append(("author", data["author"]))
                if data.get("role"):
                    a_lines.append(("role", data["role"]))
                # 用 textbox 容纳两行：用 add_paragraph
                tb = slide.shapes.add_textbox(
                    self._x(ax),
                    self._y(inner["y"]),
                    self._x(author_w - 16),
                    self._y(H),
                )
                tf = tb.text_frame
                tf.margin_left = tf.margin_right = 0
                tf.margin_top = tf.margin_bottom = 0
                tf.word_wrap = True
                from pptx.enum.text import MSO_ANCHOR

                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                for i, (kind, text) in enumerate(a_lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = text
                    run = p.runs[0] if p.runs else None
                    if run is not None:
                        if kind == "author":
                            run.font.size = Pt(_px_to_pt(20))
                            run.font.bold = True
                            run.font.color.rgb = _hex(self.theme["colors"]["text_primary"])
                        else:
                            run.font.size = Pt(_px_to_pt(14))
                            run.font.color.rgb = _hex(self.theme["colors"]["text_muted"])
                        run.font.name = "PingFang SC"
        else:
            # 大卡模式：quote 居中（PowerPoint 自动换行 + 垂直居中）
            available_h = H - (110 if data.get("author") else 30)
            self._add_textbox(
                slide,
                quote,
                inner["x"] + 32,
                inner["y"] + 30,
                W - 64,
                available_h - 30,
                font_size=q_size,
                color=self.theme["colors"]["text_primary"],
                italic=True,
                v_anchor="middle",
            )
            if data.get("author"):
                self._add_solid_rect(
                    slide,
                    inner["x"] + 32,
                    inner["y"] + H - 90,
                    48,
                    3,
                    self.theme["colors"]["accent_primary"],
                )
                self._add_textbox(
                    slide,
                    data["author"],
                    inner["x"] + 32,
                    inner["y"] + H - 70,
                    W - 64,
                    28,
                    font_size=ts["h4"],
                    color=self.theme["colors"]["text_primary"],
                    bold=True,
                )
                if data.get("role"):
                    self._add_textbox(
                        slide,
                        data["role"],
                        inner["x"] + 32,
                        inner["y"] + H - 38,
                        W - 64,
                        22,
                        font_size=ts["body"],
                        color=self.theme["colors"]["text_muted"],
                    )

    # ---------- Component: card-text ----------

    def _render_card_text(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        ts = self.theme["type_scale"]
        y = 0
        if data.get("eyebrow"):
            y += ts["eyebrow"] + 4
            self._add_textbox(
                slide,
                data["eyebrow"],
                inner["x"],
                inner["y"] + y - ts["eyebrow"],
                W,
                ts["eyebrow"] + 6,
                font_size=ts["eyebrow"],
                color=self.theme["colors"]["accent_secondary"],
                bold=True,
                letter_spacing=300,
            )
            y += 18
        if data.get("title"):
            y += ts["h3"] + 8
            self._add_textbox(
                slide,
                data["title"],
                inner["x"],
                inner["y"] + y - ts["h3"],
                W,
                int(ts["h3"] * 1.4),
                font_size=ts["h3"],
                color=self.theme["colors"]["text_primary"],
                bold=True,
            )
            y += 16
        if data.get("badges"):
            y += 22
            self._render_badges(slide, data["badges"], inner["x"], inner["y"] + y - 18)
            y += 6
        # 段落：合并成一个 textbox 多 paragraph，PowerPoint 自动换行
        paragraphs = data.get("paragraphs") or []
        if paragraphs:
            tb = slide.shapes.add_textbox(
                self._x(inner["x"]),
                self._y(inner["y"] + y),
                self._x(W),
                self._y(H - y - 8),
            )
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            tf.word_wrap = True
            for i, para in enumerate(paragraphs):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para
                p.space_after = Pt(8)
                run = p.runs[0] if p.runs else None
                if run is not None:
                    run.font.size = Pt(_px_to_pt(ts["body"]))
                    run.font.color.rgb = _hex(self.theme["colors"]["text_secondary"])
                    run.font.name = "PingFang SC"

    # ---------- Component: card-image ----------

    def _render_card_image(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        ts = self.theme["type_scale"]
        y_top = 0
        if data.get("title"):
            y_top = ts["eyebrow"] + 4
            self._add_textbox(
                slide,
                data["title"],
                inner["x"],
                inner["y"],
                W,
                ts["eyebrow"] + 6,
                font_size=ts["eyebrow"],
                color=self.theme["colors"]["accent_secondary"],
                bold=True,
                letter_spacing=300,
            )
            y_top += 16
        bottom_reserve = 30 if data.get("caption") else 8
        img_h = H - y_top - bottom_reserve
        src = data.get("src", "")
        if src and Path(src).expanduser().exists():
            try:
                slide.shapes.add_picture(
                    str(Path(src).expanduser()),
                    self._x(inner["x"]),
                    self._y(inner["y"] + y_top),
                    width=self._x(W),
                    height=self._y(img_h),
                )
            except Exception:
                self._add_image_placeholder(slide, inner, y_top, img_h, data.get("alt", "image"))
        else:
            self._add_image_placeholder(slide, inner, y_top, img_h, data.get("alt", "image"))
        if data.get("caption"):
            self._add_textbox(
                slide,
                data["caption"],
                inner["x"],
                inner["y"] + H - 22,
                W,
                20,
                font_size=ts["small"],
                color=self.theme["colors"]["text_secondary"],
            )

    def _add_image_placeholder(self, slide, inner, y_top, img_h, alt):
        ph = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self._x(inner["x"]),
            self._y(inner["y"] + y_top),
            self._x(inner["w"]),
            self._y(img_h),
        )
        self._set_solid_fill(ph, self.theme["colors"]["accent_primary"], opacity=0.10)
        self._set_line_alpha(ph, self.theme["colors"]["accent_primary"], opacity=0.4, width_pt=1)
        try:
            ph.adjustments[0] = 0.05
        except Exception:
            pass
        self._add_textbox(
            slide,
            alt,
            inner["x"],
            inner["y"] + y_top + img_h // 2 - 12,
            inner["w"],
            24,
            font_size=14,
            color=self.theme["colors"]["text_secondary"],
            align="center",
        )

    # ---------- Component: chart-bar ----------

    def _render_chart_bar(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        ts = self.theme["type_scale"]
        items = data.get("items") or []
        items_use = items[:6]
        if data.get("title"):
            self._add_textbox(
                slide,
                data["title"],
                inner["x"],
                inner["y"],
                W,
                ts["eyebrow"] + 6,
                font_size=ts["eyebrow"],
                color=self.theme["colors"]["accent_secondary"],
                bold=True,
                letter_spacing=300,
            )
            chart_y = 40
        else:
            chart_y = 8
        if not items_use:
            return
        max_val = max((it.get("value", 0) for it in items_use), default=1) or 1
        avail_h = H - chart_y - 20
        bar_h = min(48, max(24, (avail_h - (len(items_use) - 1) * 14) // len(items_use)))
        label_w = 160
        value_w = 120
        bar_w_max = W - label_w - value_w - 20
        for i, it in enumerate(items_use):
            ry = chart_y + i * (bar_h + 14)
            self._add_textbox(
                slide,
                it.get("label", ""),
                inner["x"],
                inner["y"] + ry + (bar_h - 18) // 2,
                label_w,
                22,
                font_size=ts["body"],
                color=self.theme["colors"]["text_primary"],
            )
            # 轨道
            track = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                self._x(inner["x"] + label_w),
                self._y(inner["y"] + ry),
                self._x(bar_w_max),
                self._y(bar_h),
            )
            self._set_solid_fill(track, self.theme["colors"]["card_fill"], opacity=0.08)
            track.line.fill.background()
            try:
                track.adjustments[0] = 0.15
            except Exception:
                pass
            # 填充
            v = it.get("value", 0)
            fw = int(bar_w_max * v / max_val)
            if fw > 0:
                fill_color = it.get("color") or self.theme["colors"]["accent_primary"]
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    self._x(inner["x"] + label_w),
                    self._y(inner["y"] + ry),
                    self._x(fw),
                    self._y(bar_h),
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = _hex(fill_color)
                bar.line.fill.background()
                try:
                    bar.adjustments[0] = 0.15
                except Exception:
                    pass
            unit_str = (" " + it.get("unit", "")) if it.get("unit") else ""
            self._add_textbox(
                slide,
                f"{v}{unit_str}",
                inner["x"] + label_w + fw + 12,
                inner["y"] + ry + (bar_h - 22) // 2,
                value_w,
                26,
                font_size=ts["h4"],
                color=self.theme["colors"]["text_primary"],
                bold=True,
            )

    # ---------- Component: card-compare ----------

    def _render_card_compare(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        ts = self.theme["type_scale"]
        c = self.theme["colors"]
        headers = data.get("headers") or []
        n_cols = len(headers)
        if n_cols == 0:
            return
        recommend = int(data.get("recommend", -1))
        rows = data.get("rows") or []

        y = 0
        if data.get("eyebrow"):
            self._add_textbox(
                slide,
                data["eyebrow"],
                inner["x"],
                inner["y"] + y - 4,
                W,
                ts["eyebrow"] + 6,
                font_size=ts["eyebrow"],
                color=c["accent_secondary"],
                bold=True,
                letter_spacing=300,
            )
            y += ts["eyebrow"] + 18
        if data.get("title"):
            self._add_textbox(
                slide,
                data["title"],
                inner["x"],
                inner["y"] + y - int(ts["h3"] * 0.15),
                W,
                ts["h3"] + 8,
                font_size=ts["h3"],
                color=c["text_primary"],
                bold=True,
            )
            y += ts["h3"] + 22

        table_y = y + 12
        label_w = int(W * 0.27)
        col_gap = 3
        col_w = (W - label_w - col_gap * (n_cols - 1)) // max(n_cols, 1)
        avail_h = H - table_y - 8
        row_h = max(36, min(64, int(avail_h / (len(rows) + 1.5))))
        head_h = int(row_h * 1.35)

        # 推荐列背景高亮条（贯穿所有行）
        if 0 <= recommend < n_cols:
            total_h = head_h + len(rows) * (row_h + 2) + 4
            rec_col_x = inner["x"] + label_w + recommend * (col_w + col_gap)
            hl_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                self._x(rec_col_x - 2),
                self._y(inner["y"] + table_y),
                self._x(col_w + 4),
                self._y(total_h),
            )
            self._set_solid_fill(hl_bg, c["accent_primary"], opacity=0.08)
            hl_bg.line.fill.background()
            try:
                hl_bg.adjustments[0] = 0.12
            except Exception:
                pass

        # Header 行
        for i, h in enumerate(headers):
            col_x = inner["x"] + label_w + i * (col_w + col_gap)
            hdr = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                self._x(col_x),
                self._y(inner["y"] + table_y),
                self._x(col_w),
                self._y(head_h),
            )
            if i == recommend:
                self._set_solid_fill(hdr, c["accent_primary"], opacity=1.0)
                text_color = "#ffffff" if self._is_light else "#0a0e27"
                bold = True
            else:
                self._set_solid_fill(hdr, c["text_primary"], opacity=0.08)
                text_color = c["text_secondary"]
                bold = False
            hdr.line.fill.background()
            try:
                hdr.adjustments[0] = 0.15
            except Exception:
                pass
            self._add_textbox(
                slide,
                h,
                col_x,
                inner["y"] + table_y + (head_h - ts["h4"]) // 2 - 2,
                col_w,
                ts["h4"] + 8,
                font_size=ts["h4"],
                color=text_color,
                bold=bold,
                align="center",
            )

        # 数据行
        for ri, row in enumerate(rows):
            ry = table_y + head_h + 4 + ri * (row_h + 2)
            is_hl = bool(row.get("highlight"))
            if is_hl or ri % 2 == 0:
                row_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    self._x(inner["x"]),
                    self._y(inner["y"] + ry),
                    self._x(W),
                    self._y(row_h),
                )
                self._set_solid_fill(
                    row_bg, c["accent_primary"] if is_hl else c["text_primary"], opacity=0.12 if is_hl else 0.04
                )
                row_bg.line.fill.background()
                try:
                    row_bg.adjustments[0] = 0.08
                except Exception:
                    pass
            # 行标签
            self._add_textbox(
                slide,
                row.get("label", ""),
                inner["x"],
                inner["y"] + ry + (row_h - ts["body"]) // 2 - 2,
                label_w - 8,
                ts["body"] + 6,
                font_size=ts["body"],
                color=c["text_muted"],
            )
            # 各列值
            for ci, val in enumerate(row.get("values") or []):
                col_x = inner["x"] + label_w + ci * (col_w + col_gap)
                is_rec = ci == recommend
                self._add_textbox(
                    slide,
                    str(val),
                    col_x,
                    inner["y"] + ry + (row_h - ts["body"]) // 2 - 2,
                    col_w,
                    ts["body"] + 6,
                    font_size=ts["body"],
                    color=c["text_primary"] if (is_hl or is_rec) else c["text_secondary"],
                    bold=is_hl or is_rec,
                    align="center",
                )

    def _render_card_hero(self, slide, inner: dict, data: dict) -> None:
        H, W = inner["h"], inner["w"]
        # 字号自适应（同 SVG 模板逻辑）
        ts = self.theme["type_scale"]
        if H <= 200:
            h_size, sub_size = ts["h2"], ts["h4"]
        elif H <= 320:
            h_size, sub_size = ts["h2"] + 8, ts["h3"]
        else:
            h_size, sub_size = ts["h1"], ts["h3"]
        line_h = int(h_size * 1.20)

        # 横幅模式（H<280）且无底部信息 → 中段内容整体垂直居中
        has_bottom = (H >= 280) and (data.get("meta_columns") or data.get("footer"))
        top_pad = 0
        if H < 280 and not has_bottom:
            # 估算 eyebrow + title + badges + subtitle 总高
            est = 0
            if data.get("eyebrow"):
                est += ts["eyebrow"] + 18
            if data.get("title"):
                t = data["title"]
                line_count = len(t) if isinstance(t, list) else 1
                est += h_size + (line_h * (line_count - 1))
            if data.get("badges"):
                est += 22 + 14
            if data.get("subtitle"):
                est += sub_size + 18
            top_pad = max(0, (H - est) // 2)

        # 装饰大字 deco_text（仅 H>=320，半透明大字放右下，先画在底层）
        if data.get("deco_text") and H >= 320:
            dt_size = int(H * 0.32)
            self._add_textbox(
                slide,
                data["deco_text"],
                inner["x"],
                inner["y"] + H - dt_size - 16,
                W,
                dt_size + 24,
                font_size=dt_size,
                color=self.theme["colors"]["accent_primary"],
                bold=True,
                align="right",
                transparency=0.93,
                letter_spacing=600,
            )

        y = top_pad
        # eyebrow
        if data.get("eyebrow"):
            y += ts["eyebrow"] + 4
            self._add_textbox(
                slide,
                data["eyebrow"],
                inner["x"],
                inner["y"] + y - ts["eyebrow"],
                W,
                ts["eyebrow"] + 6,
                font_size=ts["eyebrow"],
                color=self.theme["colors"]["accent_secondary"],
                bold=True,
                letter_spacing=300,
            )
            y += 18

        # title (单行 string 或 多行 list)
        title = data.get("title")
        if title:
            if isinstance(title, str):
                title_lines = [title]
            elif isinstance(title, list):
                title_lines = [str(line) for line in title]
            else:
                title_lines = [str(title)]
            for line in title_lines:
                y += h_size
                self._add_textbox(
                    slide,
                    line,
                    inner["x"],
                    inner["y"] + y - h_size,
                    W,
                    line_h,
                    font_size=h_size,
                    color=self.theme["colors"]["text_primary"],
                    bold=True,
                )
                y += line_h - h_size

        # badges
        if data.get("badges"):
            y += 22
            self._render_badges(slide, data["badges"], inner["x"], inner["y"] + y - 18)
            y += 14

        # subtitle
        if data.get("subtitle"):
            y += sub_size + 18
            self._add_textbox(
                slide,
                data["subtitle"],
                inner["x"],
                inner["y"] + y - sub_size,
                W,
                int(sub_size * 1.4),
                font_size=sub_size,
                color=self.theme["colors"]["text_secondary"],
            )

        # 底部装饰条 + footer/meta_columns（仅 H>=280）
        if H >= 280:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                self._x(inner["x"]),
                self._y(inner["y"] + H - 80),
                self._x(120),
                self._y(4),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = _hex(self.theme["colors"]["accent_primary"])
            bar.line.fill.background()

            if data.get("meta_columns"):
                cols = data["meta_columns"]
                col_w = W // len(cols)
                for i, c in enumerate(cols):
                    cx = inner["x"] + i * col_w
                    self._add_textbox(
                        slide,
                        c.get("label", ""),
                        cx,
                        inner["y"] + H - 50,
                        col_w,
                        18,
                        font_size=12,
                        color=self.theme["colors"]["text_muted"],
                        letter_spacing=200,
                    )
                    self._add_textbox(
                        slide,
                        c.get("value", ""),
                        cx,
                        inner["y"] + H - 28,
                        col_w,
                        22,
                        font_size=18,
                        color=self.theme["colors"]["text_primary"],
                        bold=True,
                    )
            elif data.get("footer"):
                self._add_textbox(
                    slide,
                    data["footer"],
                    inner["x"],
                    inner["y"] + H - 26,
                    W,
                    int(ts["body"] * 1.4),
                    font_size=ts["body"],
                    color=self.theme["colors"]["text_secondary"],
                )


def render_pptx(ws: Path, theme_name: str | None = None, out_path: Path | None = None) -> Path:
    """工作区入口：读 layout.json，渲染到 deck.pptx。"""
    layout = json.loads((ws / "layout.json").read_text(encoding="utf-8"))
    theme = theme_name or layout.get("theme", "bento-tech")
    out = Path(out_path) if out_path else (ws / "deck.pptx")
    NativeRenderer(theme).render_deck(layout, out)
    return out
